"""
Fix GFSA slide formatting and add Agenda slide to Aibel presentation.
"""
import os
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

dst = r'C:\Users\SondreFalch\OneDrive - Sifab AS\Dokumenter - Felles\A MIDLERTIDIG\Presentation Sifab for Aibel March 2026.pptx'
prs = Presentation(dst)

print(f"Loaded: {len(prs.slides)} slides")

# ====================================================================
# STEP 1: Extract logo image from an existing slide to reuse
# ====================================================================
# The logo appears on most slides at position ~L:11.67 T:0.24 W:1.51 H:0.29
logo_blob = None
logo_content_type = None
for shape in prs.slides[2].shapes:  # Q.Sonic slide
    if shape.shape_type == 13:  # Picture
        if abs(shape.left/914400 - 11.67) < 0.5 and abs(shape.top/914400 - 0.24) < 0.5:
            try:
                logo_blob = shape.image.blob
                logo_content_type = shape.image.content_type
                print(f"Found logo: {logo_content_type}, {len(logo_blob)} bytes")
            except:
                pass
            break

# Also grab the small Sifab logo from FH-Sonic slide (top-left, small)
sifab_logo_blob = None
for shape in prs.slides[5].shapes:  # FH-Sonic slide
    if shape.shape_type == 13:
        if shape.width/914400 < 1.5 and shape.top/914400 < 0.5 and shape.left/914400 < 0.5:
            try:
                sifab_logo_blob = shape.image.blob
                print(f"Found Sifab logo: {shape.image.content_type}, {len(sifab_logo_blob)} bytes")
            except:
                pass
            break

# ====================================================================
# STEP 2: Delete current GFSA slide (index 9) and Q&A slide (index 15)
# ====================================================================
# Delete in reverse order
for idx in [15, 9]:
    rId = prs.slides._sldIdLst[idx].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[idx]

print(f"After deleting old GFSA+QA: {len(prs.slides)} slides")

# Save intermediate to avoid duplicate slide name issues
tmp = dst.replace('.pptx', '_tmp.pptx')
prs.save(tmp)
prs = Presentation(tmp)
print(f"Reopened: {len(prs.slides)} slides")

# Current order (14 slides):
# 0-Title, 1-Products, 2-QSonic, 3-TZN, 4-AntiScale, 5-FHSonic,
# 6-Ultraflux, 7-VCone, 8-Emerson, 9-SVP, 10-TankGauging,
# 11-Servo, 12-CIU, 13-SmartRadar

# ====================================================================
# STEP 3: Reference slide for layout - use Emerson (idx 8) as template
# ====================================================================
# Emerson layout analysis:
# - Title: L:0.83 T:0.83 W:6.64 H:1.87, ~40pt
# - Body:  L:0.83 T:2.70 W:6.64 H:4.12, ~20pt
# - Right bg: L:8.30 T:0.00 W:5.04 H:7.50 (colored panel)
# - Images right side: L:8.97-9.70, spread vertically
# - Logo: L:11.67 T:0.24

layout = prs.slide_layouts[5]  # Use a blank layout to avoid placeholder conflicts

# ====================================================================
# STEP 4: Create AGENDA slide (will become slide 2)
# ====================================================================
agenda_slide = prs.slides.add_slide(layout)

# Title
txBox = agenda_slide.shapes.add_textbox(
    Inches(0.83), Inches(0.60), Inches(11.5), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Agenda"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(0x00, 0x64, 0x00)

# Agenda items - two columns
agenda_items_left = [
    ("1.", "Q.Sonic Max 6 Ultrasonic Gas Flowmeter"),
    ("2.", "Faure Herman Anti Scaling and Downsizing"),
    ("3.", "Faure Herman FHSonic Fiscal Ultrasonic Flowmeter"),
    ("4.", "Faure Herman Ultraflux Clamp On Flowmeters"),
    ("5.", "McCrometer V-Cone Diff. Pressure Flowmeter"),
]
agenda_items_right = [
    ("6.", "Emerson Coriolis Flowmeters"),
    ("7.", "GFSA Filters and Strainers"),
    ("8.", "Honeywell Enraf Small Volume Provers"),
    ("9.", "Honeywell Enraf Tank Gauging / Terminal Management"),
]

# Left column
txBox_left = agenda_slide.shapes.add_textbox(
    Inches(0.83), Inches(1.80), Inches(5.5), Inches(5.0))
tf_left = txBox_left.text_frame
tf_left.word_wrap = True

for i, (num, item) in enumerate(agenda_items_left):
    p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
    p.space_before = Pt(8)
    p.space_after = Pt(8)

    run_num = p.add_run()
    run_num.text = f"  {num}  "
    run_num.font.size = Pt(22)
    run_num.font.bold = True
    run_num.font.color.rgb = RGBColor(0x00, 0x64, 0x00)

    run_text = p.add_run()
    run_text.text = item
    run_text.font.size = Pt(20)
    run_text.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# Right column
txBox_right = agenda_slide.shapes.add_textbox(
    Inches(6.80), Inches(1.80), Inches(5.5), Inches(5.0))
tf_right = txBox_right.text_frame
tf_right.word_wrap = True

for i, (num, item) in enumerate(agenda_items_right):
    p = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
    p.space_before = Pt(8)
    p.space_after = Pt(8)

    run_num = p.add_run()
    run_num.text = f"  {num}  "
    run_num.font.size = Pt(22)
    run_num.font.bold = True
    run_num.font.color.rgb = RGBColor(0x00, 0x64, 0x00)

    run_text = p.add_run()
    run_text.text = item
    run_text.font.size = Pt(20)
    run_text.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# Add logo if found
if logo_blob:
    import io
    logo_stream = io.BytesIO(logo_blob)
    agenda_slide.shapes.add_picture(logo_stream, Inches(11.67), Inches(0.24), Inches(1.51))

print("  Created Agenda slide")

# ====================================================================
# STEP 5: Create new GFSA slide (matching Emerson layout)
# ====================================================================
gfsa_slide = prs.slides.add_slide(layout)

# Right panel background (matching Emerson's colored area)
from pptx.oxml import parse_xml
bg_shape = gfsa_slide.shapes.add_shape(
    1,  # MSO_SHAPE.RECTANGLE
    Inches(8.30), Inches(0), Inches(5.04), Inches(7.50)
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xF0)  # Light green-grey
bg_shape.line.fill.background()  # No border

# Title
txBox = gfsa_slide.shapes.add_textbox(
    Inches(0.83), Inches(0.83), Inches(6.64), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "7. GFSA Filters and Strainers"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0x00, 0x64, 0x00)

# Decorative line under title
line = gfsa_slide.shapes.add_shape(
    1, Inches(0.83), Inches(1.85), Inches(0.81), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x00, 0x64, 0x00)
line.line.fill.background()

# Body text
txBox2 = gfsa_slide.shapes.add_textbox(
    Inches(0.83), Inches(2.10), Inches(6.64), Inches(4.80))
tf2 = txBox2.text_frame
tf2.word_wrap = True

body_items = [
    ("GFSA Ltd. designs and manufactures filtration and process equipment "
     "for the oil & gas industry.", True),
    ("", False),
    ("Full range of filters, strainers and flame arresters \u2014 "
     "standard and bespoke designs.", False),
    ("Liquids and gas/vapour units for solid/liquid and "
     "liquid/liquid separation.", False),
    ('Size ranges from \u00bd" to 48" NB, pressure ratings up to 2500 lb.', False),
    ("Materials: Carbon steel, duplex, super duplex, 6Mo, "
     "Titanium, Hastelloy.", False),
    ("Manufactured to ASME VIII Div.1, PED, EN13445, PD5500 and NORSOK.", False),
    ("", False),
    ("Product range:", True),
    ("\u2022  Y-Type Strainers    \u2022  Basket Strainers    \u2022  Duplex Strainers", False),
    ("\u2022  Cartridge Filters    \u2022  Coalescers & Separators", False),
    ("\u2022  Flame Arresters (inline & end-of-line)", False),
    ("\u2022  Full system skid modules & pressure vessels", False),
    ("", False),
    ("SIFAB is an authorized channel partner for GFSA in Norway.", True),
]

for i, (text, bold) in enumerate(body_items):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    p.space_after = Pt(3)

# Images on right side (matching Emerson positioning)
gfsa_img_dir = r'C:\Users\SondreFalch\OneDrive - Sifab AS\Dokumenter - Felles\AGENTURER\GFSA'
images = [
    ('Picture10-Y Type Strainer Fabricated.png', Inches(8.70), Inches(0.50), Inches(3.8)),
    ('Picture17-Duplex Strainer.png', Inches(8.70), Inches(3.20), Inches(3.8)),
    ('Picture44-Skid Modules.jpg', Inches(8.70), Inches(5.30), Inches(3.8)),
]
for img_name, left, top, width in images:
    img_path = os.path.join(gfsa_img_dir, img_name)
    if os.path.exists(img_path):
        gfsa_slide.shapes.add_picture(img_path, left, top, width)
        print(f"  Added GFSA image: {img_name}")

# Add logo
if logo_blob:
    import io
    gfsa_slide.shapes.add_picture(io.BytesIO(logo_blob), Inches(11.67), Inches(0.24), Inches(1.51))

print("  Created GFSA slide (widescreen format)")

# ====================================================================
# STEP 6: Create new Q&A slide
# ====================================================================
qa_slide = prs.slides.add_slide(layout)

txBox = qa_slide.shapes.add_textbox(
    Inches(2), Inches(2.0), Inches(9), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Questions & Discussion"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0x00, 0x64, 0x00)
p.alignment = PP_ALIGN.CENTER

tf.add_paragraph().text = ""

p3 = tf.add_paragraph()
p3.text = "Thank you for your time"
p3.font.size = Pt(24)
p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p3.alignment = PP_ALIGN.CENTER

tf.add_paragraph().text = ""

p5 = tf.add_paragraph()
p5.text = "Sifab AS \u2014 Your Partner in Specialized Flow Metering"
p5.font.size = Pt(18)
p5.font.color.rgb = RGBColor(0x00, 0x64, 0x00)
p5.alignment = PP_ALIGN.CENTER

p6 = tf.add_paragraph()
p6.text = "www.sifab.no  |  Bedriftsveien 20, 4313 Sandnes"
p6.font.size = Pt(16)
p6.font.color.rgb = RGBColor(0x00, 0x64, 0x00)
p6.alignment = PP_ALIGN.CENTER

if logo_blob:
    import io
    qa_slide.shapes.add_picture(io.BytesIO(logo_blob), Inches(11.67), Inches(0.24), Inches(1.51))

print("  Created Q&A slide")

# ====================================================================
# STEP 7: Reorder slides
# ====================================================================
# Current (17 slides):
# 0-Title, 1-Products, 2-QSonic, 3-TZN, 4-AntiScale, 5-FHSonic,
# 6-Ultraflux, 7-VCone, 8-Emerson, 9-SVP, 10-TankGauging,
# 11-Servo, 12-CIU, 13-SmartRadar, 14-Agenda, 15-GFSA, 16-QA
#
# Desired order:
# Title, Agenda, Products, QSonic, TZN, AntiScale, FHSonic,
# Ultraflux, VCone, Emerson, GFSA, SVP, TankGauging, Servo, CIU,
# SmartRadar, QA
desired = [0, 14, 1, 2, 3, 4, 5, 6, 7, 8, 15, 9, 10, 11, 12, 13, 16]

sldIdLst = prs.slides._sldIdLst
elements = list(sldIdLst)
reordered = [elements[i] for i in desired]
for elem in elements:
    sldIdLst.remove(elem)
for elem in reordered:
    sldIdLst.append(elem)

# ====================================================================
# SAVE FINAL
# ====================================================================
prs.save(dst)

# Clean up temp
try:
    os.remove(tmp)
except:
    pass

# ====================================================================
# VERIFY
# ====================================================================
prs_final = Presentation(dst)
print(f"\n{'='*60}")
print(f"FINAL PRESENTATION: {len(prs_final.slides)} slides")
print(f"{'='*60}")

for i, slide in enumerate(prs_final.slides):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text[:80].strip().replace('\n', ' | ')
            print(f"  Slide {i+1:2d}: {text}")
            break
    else:
        print(f"  Slide {i+1:2d}: (images only)")
