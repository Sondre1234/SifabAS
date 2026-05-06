import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Paths
src = r'C:\Users\SondreFalch\OneDrive - Sifab AS\Dokumenter - Felles\DIVERSE SIFAB\2 MARKED\Presentation Sifab 2025.pptx'
tmp = r'C:\Users\SondreFalch\OneDrive - Sifab AS\Dokumenter - Felles\A MIDLERTIDIG\TEMP_aibel.pptx'
dst = r'C:\Users\SondreFalch\OneDrive - Sifab AS\Dokumenter - Felles\A MIDLERTIDIG\Presentation Sifab for Aibel March 2026.pptx'

# ============================================================
# STAGE 1: Delete unwanted slides, reorder, fix title, save
# ============================================================
prs = Presentation(src)
print(f"Original slides: {len(prs.slides)}")

# Delete slides (0-indexed, reverse order)
to_delete = sorted([2, 15, 16, 17, 18, 19, 20, 21, 22, 23, 24, 25], reverse=True)
for idx in to_delete:
    rId = prs.slides._sldIdLst[idx].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[idx]

print(f"After deletion: {len(prs.slides)} slides")

# Remaining (0-indexed):
# 0-Title, 1-Products, 2-FHSonic, 3-QSonic, 4-TZN, 5-AntiScale,
# 6-Ultraflux, 7-VCone, 8-SVP, 9-TankGauging, 10-Servo, 11-CIU,
# 12-SmartRadar, 13-Emerson

# Desired: Title, Products, QSonic(1), TZN(2), AntiScale(2), FHSonic(3),
#          Ultraflux(4), VCone(5), Emerson(6), SVP(8),
#          TankGauging(9), Servo(9), CIU(9), SmartRadar(9)
# (GFSA and Q&A will be added in Stage 2)
desired_order = [0, 1, 3, 4, 5, 2, 6, 7, 13, 8, 9, 10, 11, 12]

sldIdLst = prs.slides._sldIdLst
elements = list(sldIdLst)
reordered = [elements[i] for i in desired_order]
for elem in elements:
    sldIdLst.remove(elem)
for elem in reordered:
    sldIdLst.append(elem)

# Fix title slide
slide0 = prs.slides[0]
for shape in slide0.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if "Fagdag" in para.text or "Presentasjon" in para.text:
                # Combine all runs, replace, put in first run, clear rest
                combined = "".join(run.text for run in para.runs)
                combined = combined.replace("Presentasjon Fagdag 2025", "Presentation for Aibel - March 2026")
                combined = combined.replace("Presentasjon", "Presentation for Aibel - March 2026")
                combined = combined.replace("Fagdag 2025", "").strip()
                if para.runs:
                    para.runs[0].text = combined
                    for run in para.runs[1:]:
                        run.text = ""
                print(f"  Title: {combined}")

# Print slide list
print("\nStage 1 slides:")
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            print(f"  {i+1:2d}: {shape.text_frame.text[:70].strip()}")
            break

prs.save(tmp)
print(f"\nSaved stage 1 to: {tmp}")

# ============================================================
# STAGE 2: Reopen, add GFSA + Q&A, reorder, save final
# ============================================================
prs = Presentation(tmp)
print(f"\nStage 2 opened: {len(prs.slides)} slides")

# Current order after Stage 1:
# 0-Title, 1-Products, 2-QSonic, 3-TZN, 4-AntiScale, 5-FHSonic,
# 6-Ultraflux, 7-VCone, 8-Emerson, 9-SVP, 10-TankGauging,
# 11-Servo, 12-CIU, 13-SmartRadar

layout = prs.slide_layouts[1]

# --- GFSA slide ---
gfsa_slide = prs.slides.add_slide(layout)
for ph in list(gfsa_slide.placeholders):
    ph._element.getparent().remove(ph._element)

txBox = gfsa_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "7. GFSA Filters and Strainers"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x00, 0x50, 0x00)

txBox2 = gfsa_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(4.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
bullets = [
    "GFSA Ltd. designs and manufactures filtration and process equipment for the oil & gas industry.",
    "Full range of filters, strainers and flame arresters \u2014 standard and bespoke.",
    "Liquids and gas/vapour units for solid/liquid and liquid/liquid separation.",
    'Size ranges from \u00bd" to 48" nb, pressure ratings up to 2500 lb.',
    "Materials: Carbon steel, duplex, super duplex, 6Mo, Titanium, Hastelloy.",
    "Manufactured to ASME VIII Div.1, PED, EN13445, PD5500 and NORSOK.",
    "",
    "Product range:",
    "   \u2022 Y-Type Strainers (cast & fabricated)",
    "   \u2022 Inline Basket Strainers",
    "   \u2022 Duplex Strainers",
    "   \u2022 Temporary Strainers",
    "   \u2022 Cartridge Filters",
    "   \u2022 Coalescers & Separators",
    "   \u2022 Flame Arresters (inline & end-of-line)",
    "   \u2022 Full system skid modules",
    "",
    "SIFAB is an authorized channel partner for GFSA in Norway.",
]
for i, bullet in enumerate(bullets):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = bullet
    p.font.size = Pt(11)
    p.space_after = Pt(2)

gfsa_img_dir = r'C:\Users\SondreFalch\OneDrive - Sifab AS\Dokumenter - Felles\AGENTURER\GFSA'
for img_name, left, top, width in [
    ('Picture10-Y Type Strainer Fabricated.png', Inches(5.2), Inches(1.2), Inches(2.2)),
    ('Picture17-Duplex Strainer.png', Inches(5.2), Inches(3.5), Inches(2.2)),
    ('Picture44-Skid Modules.jpg', Inches(5.2), Inches(5.8), Inches(2.2)),
]:
    img_path = os.path.join(gfsa_img_dir, img_name)
    if os.path.exists(img_path):
        gfsa_slide.shapes.add_picture(img_path, left, top, width)
        print(f"  Added GFSA image: {img_name}")

# --- Q&A slide ---
qa_slide = prs.slides.add_slide(layout)
for ph in list(qa_slide.placeholders):
    ph._element.getparent().remove(ph._element)

txBox = qa_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(5.5), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Questions & Discussion"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0x00, 0x50, 0x00)
p.alignment = PP_ALIGN.CENTER
tf.add_paragraph().text = ""
p3 = tf.add_paragraph()
p3.text = "Thank you for your time"
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p3.alignment = PP_ALIGN.CENTER
tf.add_paragraph().text = ""
p5 = tf.add_paragraph()
p5.text = "Sifab AS \u2014 Your Partner in Specialized Flow Metering"
p5.font.size = Pt(14)
p5.font.color.rgb = RGBColor(0x00, 0x50, 0x00)
p5.alignment = PP_ALIGN.CENTER
p6 = tf.add_paragraph()
p6.text = "www.sifab.no"
p6.font.size = Pt(14)
p6.font.color.rgb = RGBColor(0x00, 0x50, 0x00)
p6.alignment = PP_ALIGN.CENTER

print(f"After adding slides: {len(prs.slides)} slides")

# Reorder: insert GFSA (idx 14) after Emerson (idx 8), Q&A (idx 15) stays at end
# Current: 0-13 (original), 14 (GFSA), 15 (Q&A)
# Desired: 0,1,2,3,4,5,6,7,8,14,9,10,11,12,13,15
desired_order_2 = [0, 1, 2, 3, 4, 5, 6, 7, 8, 14, 9, 10, 11, 12, 13, 15]

sldIdLst = prs.slides._sldIdLst
elements = list(sldIdLst)
reordered = [elements[i] for i in desired_order_2]
for elem in elements:
    sldIdLst.remove(elem)
for elem in reordered:
    sldIdLst.append(elem)

# Save final
prs.save(dst)

# Clean up temp
try:
    os.remove(tmp)
except:
    pass

# === VERIFY ===
prs_final = Presentation(dst)
print(f"\n{'='*60}")
print(f"FINAL PRESENTATION: {len(prs_final.slides)} slides")
print(f"{'='*60}")

agenda_map = {
    "Q.Sonic": "Agenda 1",
    "Helifu TZN": "Agenda 2 - Downsizing",
    "Anti Scaling": "Agenda 2 - Anti Scaling",
    "FH-Sonic": "Agenda 3",
    "SE 1790": "Agenda 4 - Clamp On",
    "V-Cone": "Agenda 5",
    "Emerson": "Agenda 6",
    "GFSA": "Agenda 7",
    "Small Volume Prover": "Agenda 8",
    "Tank Gauging": "Agenda 9a",
    "Servo Gauge": "Agenda 9b",
    "Communication Interface": "Agenda 9c",
    "SmartRadar": "Agenda 9d",
}

for i, slide in enumerate(prs_final.slides):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text[:80].strip()
            agenda = ""
            for key, val in agenda_map.items():
                if key in text:
                    agenda = f"  [{val}]"
                    break
            print(f"  Slide {i+1:2d}: {text}{agenda}")
            break
    else:
        print(f"  Slide {i+1:2d}: (no text)")
